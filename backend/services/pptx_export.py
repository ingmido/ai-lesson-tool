"""
Generates real .pptx presentations in one of 4 visual styles, mirroring the
layout patterns observed in the school's reference slide templates:

  - teal_light   : light background, teal accent bar/header (MOEYS default deck)
  - navy_dark    : dark navy background, green accent (same deck, dark variant)
  - moeys_formal : white background, deep-blue + gold, official ministry look
  - playful      : cream background with dot pattern, orange/blue, rounded shapes

Each slide in the AI-generated JSON has a "type" (cover / bullets / two_column /
cards / timeline / table / compare / section / closing) which maps to a
dedicated layout function below.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

KHMER_FONT = "Khmer OS Battambang"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.15)

STYLES = {
    "teal_light": {
        "bg": RGBColor(0xEC, 0xF4, 0xF5),
        "header_line": RGBColor(0x14, 0x8F, 0x8C),
        "accent": RGBColor(0x14, 0x8F, 0x8C),
        "accent_soft": RGBColor(0xDD, 0xF0, 0xEF),
        "text": RGBColor(0x1B, 0x22, 0x2E),
        "subtext": RGBColor(0x5B, 0x6B, 0x7A),
        "card_bg": RGBColor(0xFF, 0xFF, 0xFF),
        "cover_bg": RGBColor(0x0D, 0x1B, 0x2A),
        "cover_text": RGBColor(0xFF, 0xFF, 0xFF),
        "cover_accent": RGBColor(0x2E, 0xC9, 0x9A),
    },
    "navy_dark": {
        "bg": RGBColor(0x0B, 0x13, 0x22),
        "header_line": RGBColor(0x22, 0xC5, 0x8B),
        "accent": RGBColor(0x22, 0xC5, 0x8B),
        "accent_soft": RGBColor(0x14, 0x2A, 0x24),
        "text": RGBColor(0xF3, 0xF6, 0xFA),
        "subtext": RGBColor(0xA9, 0xB4, 0xC2),
        "card_bg": RGBColor(0x14, 0x1F, 0x33),
        "cover_bg": RGBColor(0x07, 0x0D, 0x18),
        "cover_text": RGBColor(0xFF, 0xFF, 0xFF),
        "cover_accent": RGBColor(0x22, 0xC5, 0x8B),
    },
    "moeys_formal": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "header_line": RGBColor(0x1B, 0x3A, 0x8F),
        "accent": RGBColor(0x1B, 0x3A, 0x8F),
        "accent_soft": RGBColor(0xEC, 0xF0, 0xFB),
        "gold": RGBColor(0xD9, 0xA5, 0x2C),
        "text": RGBColor(0x1B, 0x2A, 0x4A),
        "subtext": RGBColor(0x5B, 0x6B, 0x7A),
        "card_bg": RGBColor(0xF4, 0xF7, 0xFC),
        "cover_bg": RGBColor(0x14, 0x2C, 0x70),
        "cover_text": RGBColor(0xFF, 0xFF, 0xFF),
        "cover_accent": RGBColor(0xD9, 0xA5, 0x2C),
    },
    "playful": {
        "bg": RGBColor(0xFD, 0xF6, 0xEC),
        "header_line": RGBColor(0xF2, 0x6B, 0x38),
        "accent": RGBColor(0xF2, 0x6B, 0x38),
        "accent_soft": RGBColor(0xFF, 0xE7, 0xD6),
        "accent2": RGBColor(0x3E, 0x8E, 0xDE),
        "text": RGBColor(0x33, 0x2B, 0x24),
        "subtext": RGBColor(0x8A, 0x7A, 0x6C),
        "card_bg": RGBColor(0xFF, 0xFF, 0xFF),
        "cover_bg": RGBColor(0xFF, 0xE9, 0xD6),
        "cover_text": RGBColor(0x33, 0x2B, 0x24),
        "cover_accent": RGBColor(0xF2, 0x6B, 0x38),
    },
}

DEFAULT_STYLE = "teal_light"


def _style(name):
    return STYLES.get(name, STYLES[DEFAULT_STYLE])


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False


def _rect(slide, x, y, w, h, fill=None, shape_type=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    _no_line(shp)
    return shp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          line_spacing=1.15, word_wrap=True):
    """runs: list of paragraphs, each a list of (text, size, color, bold) tuples."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    for i, para_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for (text, size, color, bold) in para_runs:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = KHMER_FONT
    return box


def _header(slide, sty, footer_left, section_label):
    _text(slide, MARGIN, Inches(0.28), Inches(6), Inches(0.35),
          [[(footer_left or "", 11, sty["subtext"], False)]])
    _text(slide, SLIDE_W - Inches(6.5) - MARGIN, Inches(0.28), Inches(6.5), Inches(0.35),
          [[(section_label or "", 11, sty["subtext"], False)]], align=PP_ALIGN.RIGHT)
    line = _rect(slide, MARGIN, Inches(0.68), SLIDE_W - 2 * MARGIN, Pt(1.5), fill=sty["header_line"])


def _heading(slide, sty, text, y=Inches(0.9)):
    _rect(slide, MARGIN, y + Inches(0.06), Inches(0.09), Inches(0.42), fill=sty["accent"])
    _text(slide, MARGIN + Inches(0.25), y, Inches(11.5), Inches(0.6),
          [[(text, 26, sty["text"], True)]])


def _new_slide(prs, sty):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, sty["bg"])
    return slide


# ---------- Slide type renderers ----------

def render_cover(prs, sty, deck, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, sty["cover_bg"])
    _text(slide, MARGIN, Inches(0.3), Inches(6), Inches(0.35),
          [[(deck.get("footer_left", ""), 11, sty["cover_text"], False)]])
    _text(slide, SLIDE_W - Inches(6.5) - MARGIN, Inches(0.3), Inches(6.5), Inches(0.35),
          [[(s.get("section_label", ""), 11, sty["cover_accent"], True)]], align=PP_ALIGN.RIGHT)
    _rect(slide, MARGIN, Inches(0.72), SLIDE_W - 2 * MARGIN, Pt(1.5), fill=sty["cover_accent"])

    _text(slide, Inches(1.0), Inches(2.9), SLIDE_W - Inches(2.0), Inches(1.6),
          [[(s.get("heading", ""), 40, sty["cover_text"], True)]], align=PP_ALIGN.CENTER)
    if s.get("subheading"):
        _text(slide, Inches(1.5), Inches(4.15), SLIDE_W - Inches(3.0), Inches(0.6),
              [[(s.get("subheading", ""), 16, sty["cover_accent"], False)]], align=PP_ALIGN.CENTER)


def render_section(prs, sty, deck, s):
    slide = _new_slide(prs, sty)
    _rect(slide, SLIDE_W / 2 - Inches(0.6), Inches(3.05), Inches(1.2), Pt(3), fill=sty["accent"])
    _text(slide, Inches(1.2), Inches(3.3), SLIDE_W - Inches(2.4), Inches(0.9),
          [[(s.get("heading", ""), 30, sty["text"], True)]], align=PP_ALIGN.CENTER)
    if s.get("subheading"):
        _text(slide, Inches(1.6), Inches(4.2), SLIDE_W - Inches(3.2), Inches(0.6),
              [[(s.get("subheading", ""), 15, sty["subtext"], False)]], align=PP_ALIGN.CENTER)


def render_bullets(prs, sty, deck, s):
    slide = _new_slide(prs, sty)
    _header(slide, sty, deck.get("footer_left", ""), s.get("section_label", ""))
    _heading(slide, sty, s.get("heading", ""))

    items = s.get("items", [])
    y = CONTENT_TOP + Inches(0.5)
    row_h = min(Inches(1.0), (SLIDE_H - y - Inches(0.5)) / max(len(items), 1))
    for item in items:
        _rect(slide, MARGIN, y + Inches(0.08), Inches(0.32), Inches(0.32),
              fill=sty["accent"], shape_type=MSO_SHAPE.OVAL)
        _text(slide, MARGIN + Inches(0.55), y, Inches(11.6), row_h,
              [[(item.get("title", "") + "  ", 15, sty["accent"], True),
                (item.get("desc", ""), 14, sty["text"], False)]])
        y += row_h


def render_two_column(prs, sty, deck, s):
    slide = _new_slide(prs, sty)
    _header(slide, sty, deck.get("footer_left", ""), s.get("section_label", ""))
    _heading(slide, sty, s.get("heading", ""))

    left_w = Inches(6.6)
    y = CONTENT_TOP + Inches(0.5)
    paras = [[(p, 14, sty["text"], False)] for p in s.get("paragraphs", [])]
    _text(slide, MARGIN, y, left_w, Inches(4.5), paras, line_spacing=1.3)

    right_x = MARGIN + left_w + Inches(0.4)
    right_w = SLIDE_W - right_x - MARGIN
    card = _rect(slide, right_x, y, right_w, Inches(3.6), fill=sty["card_bg"],
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    if s.get("stat_number"):
        _text(slide, right_x, y + Inches(1.1), right_w, Inches(1.0),
              [[(s.get("stat_number", ""), 36, sty["accent"], True)]], align=PP_ALIGN.CENTER)
        _text(slide, right_x, y + Inches(2.0), right_w, Inches(0.5),
              [[(s.get("stat_label", ""), 13, sty["subtext"], False)]], align=PP_ALIGN.CENTER)
    elif s.get("image_caption"):
        _text(slide, right_x, y + Inches(1.5), right_w, Inches(0.6),
              [[(s.get("image_caption", ""), 13, sty["subtext"], False)]], align=PP_ALIGN.CENTER)


def render_cards(prs, sty, deck, s):
    slide = _new_slide(prs, sty)
    _header(slide, sty, deck.get("footer_left", ""), s.get("section_label", ""))
    _heading(slide, sty, s.get("heading", ""))

    cards = s.get("cards", [])[:5] or [{}]
    n = len(cards)
    gap = Inches(0.25)
    total_w = SLIDE_W - 2 * MARGIN
    card_w = (total_w - gap * (n - 1)) / n
    y = CONTENT_TOP + Inches(0.7)
    card_h = Inches(2.6)

    for i, c in enumerate(cards):
        x = MARGIN + i * (card_w + gap)
        _rect(slide, x, y, card_w, card_h, fill=sty["card_bg"], shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(slide, x + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), Inches(0.6),
              [[(c.get("icon", "•"), 22, sty["accent"], False)]])
        _text(slide, x + Inches(0.2), y + Inches(0.9), card_w - Inches(0.4), Inches(0.7),
              [[(c.get("title", ""), 13, sty["text"], True)]], line_spacing=1.1)
        _text(slide, x + Inches(0.2), y + Inches(1.55), card_w - Inches(0.4), Inches(0.9),
              [[(c.get("desc", ""), 11, sty["subtext"], False)]], line_spacing=1.15)


def render_timeline(prs, sty, deck, s):
    slide = _new_slide(prs, sty)
    _header(slide, sty, deck.get("footer_left", ""), s.get("section_label", ""))
    _heading(slide, sty, s.get("heading", ""))

    stages = s.get("stages", [])[:5] or [{}]
    n = len(stages)
    line_y = Inches(4.3)
    _rect(slide, MARGIN, line_y, SLIDE_W - 2 * MARGIN, Pt(2), fill=sty["accent"])

    total_w = SLIDE_W - 2 * MARGIN
    slot_w = total_w / n
    card_w = slot_w - Inches(0.3)
    card_h = Inches(1.5)

    for i, st in enumerate(stages):
        cx = MARGIN + i * slot_w + slot_w / 2
        _rect(slide, cx - Inches(0.11), line_y - Inches(0.11), Inches(0.22), Inches(0.22),
              fill=sty["accent"], shape_type=MSO_SHAPE.OVAL)

        card_x = MARGIN + i * slot_w + Inches(0.15)
        card_y = line_y + Inches(0.35) if i % 2 == 0 else line_y - Inches(0.35) - card_h
        _rect(slide, card_x, card_y, card_w, card_h, fill=sty["card_bg"],
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(slide, card_x + Inches(0.15), card_y + Inches(0.1), card_w - Inches(0.3), Inches(0.35),
              [[(st.get("title", ""), 12, sty["accent"], True)]])
        _text(slide, card_x + Inches(0.15), card_y + Inches(0.42), card_w - Inches(0.3), Inches(0.3),
              [[(st.get("subtitle", ""), 11, sty["text"], True)]])
        _text(slide, card_x + Inches(0.15), card_y + Inches(0.72), card_w - Inches(0.3), Inches(0.7),
              [[(st.get("desc", ""), 9.5, sty["subtext"], False)]], line_spacing=1.1)


def render_table(prs, sty, deck, s):
    slide = _new_slide(prs, sty)
    _header(slide, sty, deck.get("footer_left", ""), s.get("section_label", ""))
    _heading(slide, sty, s.get("heading", ""))

    columns = s.get("columns", [])
    rows = s.get("rows", [])
    n_rows = len(rows) + 1
    n_cols = max(len(columns), 1)

    y = CONTENT_TOP + Inches(0.5)
    w = SLIDE_W - 2 * MARGIN
    h = min(Inches(4.6), Inches(0.6) * n_rows)

    gtable = slide.shapes.add_table(n_rows, n_cols, MARGIN, y, w, h).table
    for c, col_name in enumerate(columns):
        cell = gtable.cell(0, c)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = sty["accent"]
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.runs[0].font.name = KHMER_FONT

    for r, row in enumerate(rows, start=1):
        for c in range(n_cols):
            val = row[c] if c < len(row) else ""
            cell = gtable.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = sty["card_bg"]
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(11)
            p.runs[0].font.color.rgb = sty["text"]
            p.runs[0].font.name = KHMER_FONT


def render_compare(prs, sty, deck, s):
    slide = _new_slide(prs, sty)
    _header(slide, sty, deck.get("footer_left", ""), s.get("section_label", ""))
    _heading(slide, sty, s.get("heading", ""))

    y = CONTENT_TOP + Inches(0.6)
    gap = Inches(0.4)
    w = (SLIDE_W - 2 * MARGIN - gap) / 2
    h = Inches(4.2)

    for i, side in enumerate([s.get("left", {}), s.get("right", {})]):
        x = MARGIN + i * (w + gap)
        _rect(slide, x, y, w, h, fill=sty["card_bg"], shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(slide, x + Inches(0.3), y + Inches(0.25), w - Inches(0.6), Inches(0.5),
              [[(side.get("icon", "•") + "  ", 16, sty["accent"], False),
                (side.get("title", ""), 15, sty["text"], True)]])
        by = y + Inches(0.9)
        for b in side.get("bullets", []):
            _text(slide, x + Inches(0.3), by, w - Inches(0.6), Inches(0.4),
                  [[("• ", 12, sty["accent"], True), (b, 12, sty["text"], False)]], line_spacing=1.2)
            by += Inches(0.45)


def render_closing(prs, sty, deck, s):
    slide = _new_slide(prs, sty)
    _header(slide, sty, deck.get("footer_left", ""), s.get("section_label", ""))
    _text(slide, Inches(1.0), Inches(1.6), SLIDE_W - Inches(2.0), Inches(0.9),
          [[(s.get("heading", ""), 30, sty["text"], True)]], align=PP_ALIGN.CENTER)

    y = Inches(2.8)
    for i, q in enumerate(s.get("questions", []), start=1):
        _text(slide, Inches(1.5), y, SLIDE_W - Inches(3.0), Inches(0.45),
              [[(f"{i}. ", 14, sty["accent"], True), (q, 14, sty["text"], False)]], align=PP_ALIGN.CENTER)
        y += Inches(0.5)

    if s.get("contact"):
        _text(slide, Inches(1.5), y + Inches(0.4), SLIDE_W - Inches(3.0), Inches(0.4),
              [[(s.get("contact", ""), 13, sty["accent"], True)]], align=PP_ALIGN.CENTER)


RENDERERS = {
    "cover": render_cover,
    "section": render_section,
    "bullets": render_bullets,
    "two_column": render_two_column,
    "cards": render_cards,
    "timeline": render_timeline,
    "table": render_table,
    "compare": render_compare,
    "closing": render_closing,
}


def build_pptx(data, style, output_path):
    sty = _style(style)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides_data = data.get("slides", [])
    if not slides_data or slides_data[0].get("type") != "cover":
        slides_data = [{"type": "cover", "heading": data.get("title", ""),
                         "subheading": data.get("subtitle", "")}] + slides_data

    for s in slides_data:
        renderer = RENDERERS.get(s.get("type"), render_bullets)
        renderer(prs, sty, data, s)

    prs.save(output_path)
    return output_path
